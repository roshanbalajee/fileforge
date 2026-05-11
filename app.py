import os
import uuid
import re
import csv
import json
import hashlib
import secrets
import urllib.request
import urllib.error
from io import BytesIO
from html.parser import HTMLParser
from flask import Flask, render_template, request, send_file, redirect, url_for, flash, session
from werkzeug.utils import secure_filename
from PIL import Image
from PyPDF2 import PdfMerger, PdfReader, PdfWriter
import time
from pdf2image import convert_from_path
import zipfile
import fitz # PyMuPDF

try:
    from docx import Document
except ImportError:
    Document = None

try:
    from openpyxl import Workbook, load_workbook
except ImportError:
    Workbook = None
    load_workbook = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from authlib.integrations.flask_client import OAuth
except ImportError:
    OAuth = None

app = Flask(__name__)
app.secret_key = os.environ.get("FLASK_SECRET_KEY", "fileforge_super_secret_key")

# Configuration
UPLOAD_FOLDER = 'uploads'
OUTPUT_FOLDER = 'outputs'
ALLOWED_IMAGE_EXTENSIONS = {'png', 'jpg', 'jpeg', 'webp'}
ALLOWED_PDF_EXTENSIONS = {'pdf'}
ALLOWED_DOCUMENT_EXTENSIONS = {'docx', 'txt'}
ALLOWED_SPREADSHEET_EXTENSIONS = {'xlsx', 'csv'}
ALLOWED_PRESENTATION_EXTENSIONS = {'pptx'}
POPPLER_PATH = r"C:\poppler\Library\bin" # Manual path for Poppler

app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
app.config['OUTPUT_FOLDER'] = OUTPUT_FOLDER
app.config['MAX_CONTENT_LENGTH'] = 32 * 1024 * 1024
app.config['GOOGLE_CLIENT_ID'] = os.environ.get('GOOGLE_CLIENT_ID')
app.config['GOOGLE_CLIENT_SECRET'] = os.environ.get('GOOGLE_CLIENT_SECRET')
app.config['ADSENSE_CLIENT_ID'] = os.environ.get('ADSENSE_CLIENT_ID')
app.config['ADSENSE_TOP_SLOT'] = os.environ.get('ADSENSE_TOP_SLOT')
app.config['ADSENSE_SIDE_SLOT'] = os.environ.get('ADSENSE_SIDE_SLOT')
USERS_FILE = os.path.join(os.path.dirname(__file__), 'users.json')

oauth = OAuth(app) if OAuth else None
google = None
if oauth and app.config['GOOGLE_CLIENT_ID'] and app.config['GOOGLE_CLIENT_SECRET']:
    google = oauth.register(
        name='google',
        client_id=app.config['GOOGLE_CLIENT_ID'],
        client_secret=app.config['GOOGLE_CLIENT_SECRET'],
        server_metadata_url='https://accounts.google.com/.well-known/openid-configuration',
        client_kwargs={'scope': 'openid email profile'},
    )

TOOL_CATEGORIES = [
    {
        "title": "Organize PDF",
        "tools": [
            {"name": "Merge PDF", "slug": "merge-pdf", "url": "/merge", "icon": "MG", "tone": "coral", "ready": True},
            {"name": "Split PDF", "slug": "split-pdf", "url": "/tools/split-pdf", "icon": "SP", "tone": "blue", "ready": True},
            {"name": "Remove pages", "slug": "remove-pages", "url": "/tools/remove-pages", "icon": "RM", "tone": "blue", "ready": True},
            {"name": "Extract pages", "slug": "extract-pages", "url": "/tools/extract-pages", "icon": "EX", "tone": "blue", "ready": True},
            {"name": "Organize PDF", "slug": "organize-pdf", "url": "/tools/organize-pdf", "icon": "OR", "tone": "blue", "ready": True},
            {"name": "Scan to PDF", "slug": "scan-to-pdf", "url": "/tools/scan-to-pdf", "icon": "SC", "tone": "blue", "ready": True},
        ],
    },
    {
        "title": "Optimize PDF",
        "tools": [
            {"name": "Compress PDF", "slug": "compress-pdf", "url": "/tools/compress-pdf", "icon": "CP", "tone": "blue", "ready": True},
            {"name": "Repair PDF", "slug": "repair-pdf", "url": "/tools/repair-pdf", "icon": "RP", "tone": "blue", "ready": True},
            {"name": "OCR PDF", "slug": "ocr-pdf", "url": "/tools/ocr-pdf", "icon": "OC", "tone": "blue", "ready": True},
        ],
    },
    {
        "title": "Convert to PDF",
        "tools": [
            {"name": "JPG to PDF", "slug": "jpg-to-pdf", "url": "/image-to-pdf", "icon": "JP", "tone": "gold", "ready": True},
            {"name": "WORD to PDF", "slug": "word-to-pdf", "url": "/tools/word-to-pdf", "icon": "WD", "tone": "blue", "ready": True},
            {"name": "POWERPOINT to PDF", "slug": "powerpoint-to-pdf", "url": "/tools/powerpoint-to-pdf", "icon": "PP", "tone": "blue", "ready": True},
            {"name": "EXCEL to PDF", "slug": "excel-to-pdf", "url": "/tools/excel-to-pdf", "icon": "XL", "tone": "blue", "ready": True},
            {"name": "HTML to PDF", "slug": "html-to-pdf", "url": "/tools/html-to-pdf", "icon": "HT", "tone": "blue", "ready": True},
        ],
    },
    {
        "title": "Convert from PDF",
        "tools": [
            {"name": "PDF to JPG", "slug": "pdf-to-jpg", "url": "/pdf-to-image", "icon": "PJ", "tone": "gold", "ready": True},
            {"name": "PDF to WORD", "slug": "pdf-to-word", "url": "/tools/pdf-to-word", "icon": "PW", "tone": "blue", "ready": True},
            {"name": "PDF to POWERPOINT", "slug": "pdf-to-powerpoint", "url": "/tools/pdf-to-powerpoint", "icon": "PP", "tone": "blue", "ready": True},
            {"name": "PDF to EXCEL", "slug": "pdf-to-excel", "url": "/tools/pdf-to-excel", "icon": "PX", "tone": "blue", "ready": True},
            {"name": "PDF to PDF/A", "slug": "pdf-to-pdfa", "url": "/tools/pdf-to-pdfa", "icon": "PA", "tone": "blue", "ready": True},
        ],
    },
    {
        "title": "Edit PDF",
        "tools": [
            {"name": "Rotate PDF", "slug": "rotate-pdf", "url": "/tools/rotate-pdf", "icon": "RT", "tone": "blue", "ready": True},
            {"name": "Add page numbers", "slug": "add-page-numbers", "url": "/tools/add-page-numbers", "icon": "PN", "tone": "blue", "ready": True},
            {"name": "Add watermark", "slug": "add-watermark", "url": "/tools/add-watermark", "icon": "WM", "tone": "blue", "ready": True},
            {"name": "Crop PDF", "slug": "crop-pdf", "url": "/tools/crop-pdf", "icon": "CR", "tone": "blue", "ready": True},
            {"name": "Edit PDF", "slug": "edit-pdf", "url": "/tools/edit-pdf", "icon": "ED", "tone": "blue", "ready": True},
            {"name": "PDF Forms", "slug": "pdf-forms", "url": "/tools/pdf-forms", "icon": "FM", "tone": "blue", "ready": True},
        ],
    },
    {
        "title": "PDF Security",
        "tools": [
            {"name": "Unlock PDF", "slug": "unlock-pdf", "url": "/tools/unlock-pdf", "icon": "UN", "tone": "blue", "ready": True},
            {"name": "Protect PDF", "slug": "protect-pdf", "url": "/tools/protect-pdf", "icon": "PR", "tone": "blue", "ready": True},
            {"name": "Sign PDF", "slug": "sign-pdf", "url": "/tools/sign-pdf", "icon": "SG", "tone": "blue", "ready": True},
            {"name": "Redact PDF", "slug": "redact-pdf", "url": "/tools/redact-pdf", "icon": "RD", "tone": "blue", "ready": True},
            {"name": "Compare PDF", "slug": "compare-pdf", "url": "/tools/compare-pdf", "icon": "CM", "tone": "blue", "ready": True},
        ],
    },
    {
        "title": "PDF Intelligence",
        "tools": [
            {"name": "AI Summarizer", "slug": "ai-summarizer", "url": "/tools/ai-summarizer", "icon": "AI", "tone": "blue", "ready": True},
            {"name": "Translate PDF", "slug": "translate-pdf", "url": "/tools/translate-pdf", "icon": "TR", "tone": "blue", "ready": True},
        ],
    },
]

TOOL_LOOKUP = {
    tool["slug"]: {**tool, "category": category["title"]}
    for category in TOOL_CATEGORIES
    for tool in category["tools"]
}

# Ensure directories exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(OUTPUT_FOLDER, exist_ok=True)

def allowed_file(filename, allowed_extensions):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in allowed_extensions

def cleanup_old_files():
    """Delete files older than 30 minutes in uploads and outputs."""
    now = time.time()
    for folder in [UPLOAD_FOLDER, OUTPUT_FOLDER]:
        if not os.path.exists(folder): continue
        for filename in os.listdir(folder):
            filepath = os.path.join(folder, filename)
            if os.stat(filepath).st_mtime < now - 1800:
                try:
                    os.remove(filepath)
                except:
                    pass

def pdf_to_images_fallback(pdf_path, output_folder, unique_id):
    """Fallback using PyMuPDF if pdf2image (poppler) fails."""
    doc = fitz.open(pdf_path)
    images = []
    for i, page in enumerate(doc):
        pix = page.get_pixmap()
        page_filename = f"page_{i+1}.png"
        output_file = os.path.join(output_folder, f"fallback_{unique_id}_{page_filename}")
        pix.save(output_file)
        images.append((output_file, page_filename))
    doc.close()
    return images

def unique_filename(prefix, extension):
    return f"{prefix}_{str(uuid.uuid4())[:8]}.{extension}"

def save_uploaded_file(field_name='file', allowed_extensions=ALLOWED_PDF_EXTENSIONS):
    if field_name not in request.files:
        raise ValueError('No file uploaded.')
    file = request.files[field_name]
    if file.filename == '':
        raise ValueError('No file selected.')
    if not allowed_file(file.filename, allowed_extensions):
        raise ValueError(f"Invalid file type. Allowed: {', '.join(sorted(allowed_extensions))}.")
    filename = secure_filename(file.filename)
    path = os.path.join(app.config['UPLOAD_FOLDER'], f"{str(uuid.uuid4())[:8]}_{filename}")
    file.save(path)
    return path, filename

def parse_pages(page_text, total_pages, default_all=False):
    if not page_text:
        return list(range(total_pages)) if default_all else []
    pages = set()
    for part in page_text.replace(' ', '').split(','):
        if not part:
            continue
        if '-' in part:
            start, end = part.split('-', 1)
            start = int(start)
            end = int(end)
            for page in range(start, end + 1):
                if 1 <= page <= total_pages:
                    pages.add(page - 1)
        else:
            page = int(part)
            if 1 <= page <= total_pages:
                pages.add(page - 1)
    return sorted(pages)

def save_doc(doc, prefix='processed'):
    output_filename = unique_filename(prefix, 'pdf')
    output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
    doc.save(output_path, garbage=4, deflate=True, clean=True)
    doc.close()
    return output_filename

def extract_pdf_text(pdf_path):
    doc = fitz.open(pdf_path)
    chunks = []
    for index, page in enumerate(doc, start=1):
        text = page.get_text().strip()
        if text:
            chunks.append(f"Page {index}\n{text}")
    doc.close()
    return "\n\n".join(chunks).strip()

def write_text_file(prefix, text, extension='txt'):
    output_filename = unique_filename(prefix, extension)
    output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
    with open(output_path, 'w', encoding='utf-8') as output:
        output.write(text or 'No text content was found.')
    return output_filename

def text_to_pdf(text, prefix='document'):
    output_filename = unique_filename(prefix, 'pdf')
    output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
    doc = fitz.open()
    page = doc.new_page(width=595, height=842)
    margin = 54
    cursor = margin
    for raw_line in (text or 'No readable text found.').splitlines():
        line = raw_line[:110] if raw_line else ' '
        if cursor > 790:
            page = doc.new_page(width=595, height=842)
            cursor = margin
        page.insert_text((margin, cursor), line, fontsize=11, color=(0.05, 0.13, 0.28))
        cursor += 16
    doc.save(output_path, garbage=4, deflate=True)
    doc.close()
    return output_filename

class SimpleHTMLTextParser(HTMLParser):
    def __init__(self):
        super().__init__()
        self.parts = []

    def handle_data(self, data):
        clean = data.strip()
        if clean:
            self.parts.append(clean)

def html_to_text(html):
    parser = SimpleHTMLTextParser()
    parser.feed(html)
    return "\n".join(parser.parts)

def docx_to_text(path):
    if Document is None:
        raise ValueError('python-docx is not installed.')
    doc = Document(path)
    return "\n".join(p.text for p in doc.paragraphs if p.text.strip())

def spreadsheet_to_text(path):
    if path.lower().endswith('.csv'):
        rows = []
        with open(path, newline='', encoding='utf-8', errors='ignore') as handle:
            for row in csv.reader(handle):
                rows.append(" | ".join(row))
        return "\n".join(rows)
    if load_workbook is None:
        raise ValueError('openpyxl is not installed.')
    workbook = load_workbook(path, data_only=True)
    lines = []
    for sheet in workbook.worksheets:
        lines.append(sheet.title)
        for row in sheet.iter_rows(values_only=True):
            values = [str(value) if value is not None else '' for value in row]
            if any(values):
                lines.append(" | ".join(values))
    return "\n".join(lines)

def presentation_to_text(path):
    if Presentation is None:
        raise ValueError('python-pptx is not installed.')
    deck = Presentation(path)
    lines = []
    for index, slide in enumerate(deck.slides, start=1):
        lines.append(f"Slide {index}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
    return "\n".join(lines)

def create_docx_from_text(text, prefix='pdf_to_word'):
    if Document is None:
        raise ValueError('python-docx is not installed.')
    output_filename = unique_filename(prefix, 'docx')
    output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
    doc = Document()
    doc.add_heading('Converted PDF Text', level=1)
    for block in (text or 'No readable text found.').split('\n\n'):
        doc.add_paragraph(block)
    doc.save(output_path)
    return output_filename

def create_xlsx_from_text(text, prefix='pdf_to_excel'):
    if Workbook is None:
        raise ValueError('openpyxl is not installed.')
    output_filename = unique_filename(prefix, 'xlsx')
    output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = 'PDF Text'
    for row_index, line in enumerate((text or 'No readable text found.').splitlines(), start=1):
        sheet.cell(row=row_index, column=1, value=line)
    workbook.save(output_path)
    return output_filename

def create_pptx_from_text(text, prefix='pdf_to_powerpoint'):
    if Presentation is None:
        raise ValueError('python-pptx is not installed.')
    output_filename = unique_filename(prefix, 'pptx')
    output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
    deck = Presentation()
    lines = [line for line in (text or 'No readable text found.').splitlines() if line.strip()]
    for chunk_start in range(0, max(len(lines), 1), 8):
        slide = deck.slides.add_slide(deck.slide_layouts[1])
        slide.shapes.title.text = 'Converted PDF'
        body = slide.placeholders[1]
        body.text = "\n".join(lines[chunk_start:chunk_start + 8]) or 'No readable text found.'
    deck.save(output_path)
    return output_filename

def simple_summary(text):
    sentences = re.split(r'(?<=[.!?])\s+', text.replace('\n', ' '))
    sentences = [sentence.strip() for sentence in sentences if len(sentence.strip()) > 30]
    if not sentences:
        return text[:1000] or 'No readable text found.'
    return "\n".join(f"- {sentence}" for sentence in sentences[:8])

def simple_translate(text, language):
    dictionaries = {
        'spanish': {'file': 'archivo', 'page': 'pagina', 'document': 'documento', 'name': 'nombre', 'date': 'fecha', 'total': 'total'},
        'hindi': {'file': 'फाइल', 'page': 'पेज', 'document': 'दस्तावेज', 'name': 'नाम', 'date': 'तारीख', 'total': 'कुल'},
        'french': {'file': 'fichier', 'page': 'page', 'document': 'document', 'name': 'nom', 'date': 'date', 'total': 'total'},
    }
    translated = text
    for source, target in dictionaries.get(language, dictionaries['spanish']).items():
        translated = re.sub(rf'\b{source}\b', target, translated, flags=re.IGNORECASE)
    return f"Local dictionary translation preview ({language.title()})\n\n{translated}"

def openai_chat(system_prompt, user_prompt):
    api_key = os.environ.get('OPENAI_API_KEY')
    if not api_key:
        return None
    payload = {
        'model': os.environ.get('OPENAI_MODEL', 'gpt-4o-mini'),
        'messages': [
            {'role': 'system', 'content': system_prompt},
            {'role': 'user', 'content': user_prompt[:18000]},
        ],
        'temperature': 0.2,
    }
    request_data = json.dumps(payload).encode('utf-8')
    api_request = urllib.request.Request(
        'https://api.openai.com/v1/chat/completions',
        data=request_data,
        headers={
            'Authorization': f'Bearer {api_key}',
            'Content-Type': 'application/json',
        },
        method='POST',
    )
    try:
        with urllib.request.urlopen(api_request, timeout=45) as response:
            data = json.loads(response.read().decode('utf-8'))
            return data['choices'][0]['message']['content'].strip()
    except (urllib.error.URLError, urllib.error.HTTPError, KeyError, IndexError, json.JSONDecodeError):
        return None

def ai_summary(text):
    ai_result = openai_chat(
        'You summarize PDF documents clearly and accurately.',
        f'Summarize this PDF text into concise bullets and include key action items if any:\n\n{text}',
    )
    return ai_result or simple_summary(text)

def ai_translate(text, language):
    ai_result = openai_chat(
        'You translate PDF text accurately while preserving formatting and meaning.',
        f'Translate this text to {language}. Keep headings and lists readable:\n\n{text}',
    )
    return ai_result or simple_translate(text, language)

def load_users():
    if not os.path.exists(USERS_FILE):
        return {}
    with open(USERS_FILE, 'r', encoding='utf-8') as handle:
        return json.load(handle)

def save_users(users):
    with open(USERS_FILE, 'w', encoding='utf-8') as handle:
        json.dump(users, handle, indent=2)

def password_hash(password, salt=None):
    salt = salt or secrets.token_hex(16)
    digest = hashlib.pbkdf2_hmac('sha256', password.encode('utf-8'), salt.encode('utf-8'), 120000).hex()
    return f"{salt}${digest}"

def verify_password(password, stored_hash):
    try:
        salt, expected = stored_hash.split('$', 1)
    except ValueError:
        return False
    return password_hash(password, salt).split('$', 1)[1] == expected

def sign_in_user(name, email):
    session['user'] = {
        'name': name,
        'email': email,
        'picture': None,
    }

@app.context_processor
def inject_site_data():
    def preview_info(filename):
        if not filename:
            return None
        ext = filename.rsplit('.', 1)[-1].lower() if '.' in filename else ''
        path = os.path.join(app.config['OUTPUT_FOLDER'], filename)
        size = os.path.getsize(path) if os.path.exists(path) else 0
        kind = 'document'
        if ext == 'pdf':
            kind = 'pdf'
        elif ext in {'png', 'jpg', 'jpeg', 'webp'}:
            kind = 'image'
        elif ext == 'txt':
            kind = 'text'
        elif ext == 'zip':
            kind = 'archive'
        return {
            'filename': filename,
            'extension': ext.upper() or 'FILE',
            'kind': kind,
            'size': format_file_size(size),
            'preview_url': url_for('preview_file', filename=filename),
            'download_url': url_for('download_file', filename=filename),
        }

    def preview_text(filename):
        if not filename or not filename.lower().endswith('.txt'):
            return ''
        path = os.path.join(app.config['OUTPUT_FOLDER'], filename)
        if not os.path.exists(path):
            return ''
        with open(path, 'r', encoding='utf-8', errors='ignore') as handle:
            return handle.read(5000)

    return {
        "tool_categories": TOOL_CATEGORIES,
        "current_user": session.get("user"),
        "adsense_client_id": app.config['ADSENSE_CLIENT_ID'],
        "adsense_top_slot": app.config['ADSENSE_TOP_SLOT'],
        "adsense_side_slot": app.config['ADSENSE_SIDE_SLOT'],
        "preview_info": preview_info,
        "preview_text": preview_text,
    }

def format_file_size(size):
    units = ['B', 'KB', 'MB', 'GB']
    value = float(size)
    for unit in units:
        if value < 1024 or unit == units[-1]:
            return f"{value:.1f} {unit}" if unit != 'B' else f"{int(value)} {unit}"
        value /= 1024

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/pricing')
def pricing():
    return render_template('pricing.html')

@app.route('/tools/<slug>', methods=['GET', 'POST'])
def tool_placeholder(slug):
    tool = TOOL_LOOKUP.get(slug)
    if not tool:
        flash('Tool not found.', 'error')
        return redirect(url_for('index'))
    if slug == 'merge-pdf':
        return redirect(url_for('pdf_merge'))
    if slug == 'jpg-to-pdf':
        return redirect(url_for('image_to_pdf'))
    if slug == 'pdf-to-jpg':
        return redirect(url_for('pdf_to_image'))
    if request.method == 'POST':
        cleanup_old_files()
        try:
            output_filename = process_tool(slug)
            return render_template('tool_placeholder.html', tool=tool, success=True, download_url=output_filename)
        except Exception as e:
            flash(f'Error processing file: {str(e)}', 'error')
            return redirect(request.url)
    return render_template('tool_placeholder.html', tool=tool)

def process_tool(slug):
    if slug == 'scan-to-pdf':
        return process_scan_to_pdf()
    if slug in {'word-to-pdf', 'excel-to-pdf', 'powerpoint-to-pdf', 'html-to-pdf'}:
        return process_to_pdf(slug)
    if slug == 'compare-pdf':
        return process_compare_pdf()

    pdf_path, _ = save_uploaded_file('file', ALLOWED_PDF_EXTENSIONS)
    if slug == 'split-pdf':
        return process_split_pdf(pdf_path)
    if slug == 'remove-pages':
        return process_remove_pages(pdf_path)
    if slug == 'extract-pages':
        return process_extract_pages(pdf_path)
    if slug == 'organize-pdf':
        return process_organize_pdf(pdf_path)
    if slug == 'compress-pdf':
        return process_compress_pdf(pdf_path)
    if slug == 'repair-pdf':
        return process_repair_pdf(pdf_path)
    if slug == 'ocr-pdf':
        return write_text_file('ocr_text', extract_pdf_text(pdf_path))
    if slug == 'pdf-to-word':
        return create_docx_from_text(extract_pdf_text(pdf_path))
    if slug == 'pdf-to-powerpoint':
        return create_pptx_from_text(extract_pdf_text(pdf_path))
    if slug == 'pdf-to-excel':
        return create_xlsx_from_text(extract_pdf_text(pdf_path))
    if slug == 'pdf-to-pdfa':
        return process_pdfa(pdf_path)
    if slug == 'rotate-pdf':
        return process_rotate_pdf(pdf_path)
    if slug == 'add-page-numbers':
        return process_add_page_numbers(pdf_path)
    if slug == 'add-watermark':
        return process_add_watermark(pdf_path)
    if slug == 'crop-pdf':
        return process_crop_pdf(pdf_path)
    if slug == 'edit-pdf':
        return process_edit_pdf(pdf_path)
    if slug == 'pdf-forms':
        return process_pdf_forms(pdf_path)
    if slug == 'unlock-pdf':
        return process_unlock_pdf(pdf_path)
    if slug == 'protect-pdf':
        return process_protect_pdf(pdf_path)
    if slug == 'sign-pdf':
        return process_sign_pdf(pdf_path)
    if slug == 'redact-pdf':
        return process_redact_pdf(pdf_path)
    if slug == 'ai-summarizer':
        return write_text_file('summary', ai_summary(extract_pdf_text(pdf_path)))
    if slug == 'translate-pdf':
        language = request.form.get('language', 'spanish')
        return write_text_file('translation', ai_translate(extract_pdf_text(pdf_path), language))
    raise ValueError('This tool is not available yet.')

def process_split_pdf(pdf_path):
    source = fitz.open(pdf_path)
    output_filename = unique_filename('split_pages', 'zip')
    output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
    with zipfile.ZipFile(output_path, 'w') as zipf:
        for index in range(source.page_count):
            doc = fitz.open()
            doc.insert_pdf(source, from_page=index, to_page=index)
            page_name = f"page_{index + 1}.pdf"
            temp_path = os.path.join(app.config['OUTPUT_FOLDER'], f"{str(uuid.uuid4())[:8]}_{page_name}")
            doc.save(temp_path)
            doc.close()
            zipf.write(temp_path, page_name)
            os.remove(temp_path)
    source.close()
    return output_filename

def process_remove_pages(pdf_path):
    doc = fitz.open(pdf_path)
    pages = parse_pages(request.form.get('pages', ''), doc.page_count)
    for page_index in sorted(pages, reverse=True):
        doc.delete_page(page_index)
    if doc.page_count == 0:
        page = doc.new_page()
        page.insert_text((72, 72), 'All selected pages were removed.', fontsize=12, color=(0.05, 0.22, 0.7))
    return save_doc(doc, 'removed_pages')

def process_extract_pages(pdf_path):
    source = fitz.open(pdf_path)
    pages = parse_pages(request.form.get('pages', ''), source.page_count, default_all=True)
    doc = fitz.open()
    for page_index in pages:
        doc.insert_pdf(source, from_page=page_index, to_page=page_index)
    source.close()
    return save_doc(doc, 'extracted_pages')

def process_organize_pdf(pdf_path):
    source = fitz.open(pdf_path)
    order = parse_pages(request.form.get('page_order', ''), source.page_count, default_all=True)
    doc = fitz.open()
    for page_index in order:
        doc.insert_pdf(source, from_page=page_index, to_page=page_index)
    source.close()
    return save_doc(doc, 'organized')

def process_compress_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    return save_doc(doc, 'compressed_pdf')

def process_repair_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    return save_doc(doc, 'repaired')

def process_pdfa(pdf_path):
    doc = fitz.open(pdf_path)
    doc.set_metadata({
        'title': 'FileForge archive PDF',
        'producer': 'FileForge',
        'creator': 'FileForge',
    })
    return save_doc(doc, 'archive_pdf')

def process_rotate_pdf(pdf_path):
    degrees = int(request.form.get('degrees', '90'))
    doc = fitz.open(pdf_path)
    for page in doc:
        page.set_rotation((page.rotation + degrees) % 360)
    return save_doc(doc, 'rotated')

def process_add_page_numbers(pdf_path):
    doc = fitz.open(pdf_path)
    for index, page in enumerate(doc, start=1):
        rect = page.rect
        page.insert_text((rect.width / 2 - 18, rect.height - 28), str(index), fontsize=11, color=(0.05, 0.22, 0.7))
    return save_doc(doc, 'numbered')

def process_add_watermark(pdf_path):
    text = request.form.get('text', 'FileForge')
    doc = fitz.open(pdf_path)
    for page in doc:
        rect = page.rect
        page.insert_text((rect.width * 0.24, rect.height * 0.52), text, fontsize=34, color=(0.35, 0.55, 0.95), fill_opacity=0.22)
    return save_doc(doc, 'watermarked')

def process_crop_pdf(pdf_path):
    margin = float(request.form.get('margin', '36'))
    doc = fitz.open(pdf_path)
    for page in doc:
        rect = page.rect
        page.set_cropbox(fitz.Rect(rect.x0 + margin, rect.y0 + margin, rect.x1 - margin, rect.y1 - margin))
    return save_doc(doc, 'cropped')

def process_edit_pdf(pdf_path):
    text = request.form.get('text', 'Edited with FileForge')
    doc = fitz.open(pdf_path)
    page = doc[0]
    page.insert_text((54, 54), text, fontsize=14, color=(0.05, 0.22, 0.7))
    return save_doc(doc, 'edited')

def process_pdf_forms(pdf_path):
    doc = fitz.open(pdf_path)
    lines = ['PDF form fields report']
    found = False
    for page_index, page in enumerate(doc, start=1):
        for widget in page.widgets() or []:
            found = True
            lines.append(f"Page {page_index}: {widget.field_name} = {widget.field_value}")
    doc.close()
    if not found:
        lines.append('No interactive form fields were found.')
    return write_text_file('pdf_forms', "\n".join(lines))

def process_unlock_pdf(pdf_path):
    password = request.form.get('password', '')
    reader = PdfReader(pdf_path)
    if reader.is_encrypted:
        if not password:
            raise ValueError('Password is required for encrypted PDFs.')
        reader.decrypt(password)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    output_filename = unique_filename('unlocked', 'pdf')
    output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
    with open(output_path, 'wb') as output:
        writer.write(output)
    return output_filename

def process_protect_pdf(pdf_path):
    password = request.form.get('password', '')
    if not password:
        raise ValueError('Enter a password to protect the PDF.')
    reader = PdfReader(pdf_path)
    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)
    writer.encrypt(password)
    output_filename = unique_filename('protected', 'pdf')
    output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
    with open(output_path, 'wb') as output:
        writer.write(output)
    return output_filename

def process_sign_pdf(pdf_path):
    signer = request.form.get('text', 'Signed with FileForge')
    doc = fitz.open(pdf_path)
    page = doc[-1]
    rect = page.rect
    page.draw_rect(fitz.Rect(rect.width - 250, rect.height - 96, rect.width - 36, rect.height - 42), color=(0.05, 0.22, 0.7), width=1)
    page.insert_text((rect.width - 236, rect.height - 64), signer, fontsize=12, color=(0.05, 0.22, 0.7))
    return save_doc(doc, 'signed')

def process_redact_pdf(pdf_path):
    text = request.form.get('text', '')
    if not text:
        raise ValueError('Enter text to redact.')
    doc = fitz.open(pdf_path)
    for page in doc:
        for area in page.search_for(text):
            page.add_redact_annot(area, fill=(0, 0, 0))
        page.apply_redactions()
    return save_doc(doc, 'redacted')

def process_compare_pdf():
    first, _ = save_uploaded_file('file', ALLOWED_PDF_EXTENSIONS)
    second, _ = save_uploaded_file('second_file', ALLOWED_PDF_EXTENSIONS)
    first_text = extract_pdf_text(first).splitlines()
    second_text = extract_pdf_text(second).splitlines()
    added = [line for line in second_text if line and line not in first_text]
    removed = [line for line in first_text if line and line not in second_text]
    report = "PDF comparison report\n\nAdded lines:\n" + "\n".join(added[:120])
    report += "\n\nRemoved lines:\n" + "\n".join(removed[:120])
    return write_text_file('pdf_compare', report)

def process_scan_to_pdf():
    files = request.files.getlist('files')
    if not files or files[0].filename == '':
        raise ValueError('Select one or more images.')
    images = []
    for file in files:
        if file and allowed_file(file.filename, ALLOWED_IMAGE_EXTENSIONS):
            image = Image.open(file.stream)
            if image.mode != 'RGB':
                image = image.convert('RGB')
            images.append(image)
    if not images:
        raise ValueError('No valid images uploaded.')
    output_filename = unique_filename('scan_to_pdf', 'pdf')
    output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
    images[0].save(output_path, save_all=True, append_images=images[1:])
    return output_filename

def process_to_pdf(slug):
    if slug == 'word-to-pdf':
        path, filename = save_uploaded_file('file', ALLOWED_DOCUMENT_EXTENSIONS)
        text = docx_to_text(path) if filename.lower().endswith('.docx') else open(path, encoding='utf-8', errors='ignore').read()
        return text_to_pdf(text, 'word_to_pdf')
    if slug == 'excel-to-pdf':
        path, _ = save_uploaded_file('file', ALLOWED_SPREADSHEET_EXTENSIONS)
        return text_to_pdf(spreadsheet_to_text(path), 'excel_to_pdf')
    if slug == 'powerpoint-to-pdf':
        path, _ = save_uploaded_file('file', ALLOWED_PRESENTATION_EXTENSIONS)
        return text_to_pdf(presentation_to_text(path), 'powerpoint_to_pdf')
    if slug == 'html-to-pdf':
        html = request.form.get('html', '')
        if not html and 'file' in request.files and request.files['file'].filename:
            file = request.files['file']
            html = file.read().decode('utf-8', errors='ignore')
        return text_to_pdf(html_to_text(html), 'html_to_pdf')
    raise ValueError('Unsupported conversion.')

@app.route('/login')
def login():
    return render_template('login.html', mode='login', google_ready=bool(google))

@app.route('/login', methods=['POST'])
def login_post():
    email = request.form.get('email', '').strip().lower()
    password = request.form.get('password', '')
    users = load_users()
    user = users.get(email)
    if not user or not verify_password(password, user.get('password_hash', '')):
        flash('Invalid email or password.', 'error')
        return redirect(url_for('login'))
    sign_in_user(user.get('name') or email.split('@')[0], email)
    flash('Logged in successfully.', 'success')
    return redirect(url_for('index'))

@app.route('/signup')
def signup():
    return render_template('login.html', mode='signup', google_ready=bool(google))

@app.route('/signup', methods=['POST'])
def signup_post():
    name = request.form.get('name', '').strip()
    email = request.form.get('email', '').strip().lower()
    password = request.form.get('password', '')
    if not name or not email or not password:
        flash('Name, email, and password are required.', 'error')
        return redirect(url_for('signup'))
    if len(password) < 6:
        flash('Password must be at least 6 characters.', 'error')
        return redirect(url_for('signup'))
    users = load_users()
    if email in users:
        flash('An account already exists for this email.', 'error')
        return redirect(url_for('login'))
    users[email] = {
        'name': name,
        'email': email,
        'password_hash': password_hash(password),
    }
    save_users(users)
    sign_in_user(name, email)
    flash('Account created successfully.', 'success')
    return redirect(url_for('index'))

@app.route('/auth/google')
def google_login():
    if not google:
        flash('Google sign-in needs GOOGLE_CLIENT_ID and GOOGLE_CLIENT_SECRET in your environment.', 'error')
        return redirect(url_for('login'))
    redirect_uri = url_for('google_callback', _external=True)
    return google.authorize_redirect(redirect_uri)

@app.route('/auth/google/callback')
def google_callback():
    if not google:
        flash('Google sign-in is not configured yet.', 'error')
        return redirect(url_for('login'))
    token = google.authorize_access_token()
    user_info = token.get('userinfo') or google.userinfo()
    sign_in_user(user_info.get('name'), user_info.get('email'))
    session['user']['picture'] = user_info.get('picture')
    flash('Signed in with Google successfully.', 'success')
    return redirect(url_for('index'))

@app.route('/logout')
def logout():
    session.pop('user', None)
    flash('Signed out successfully.', 'success')
    return redirect(url_for('index'))

@app.route('/compress', methods=['GET', 'POST'])
def image_compress():
    if request.method == 'POST':
        cleanup_old_files()
        if 'file' not in request.files:
            flash('No file part', 'error')
            return redirect(request.url)
        
        file = request.files['file']
        quality_val = request.form.get('quality', 'medium')
        
        if file.filename == '':
            flash('No selected file', 'error')
            return redirect(request.url)
        
        if file and allowed_file(file.filename, ALLOWED_IMAGE_EXTENSIONS):
            filename = secure_filename(file.filename)
            unique_id = str(uuid.uuid4())[:8]
            input_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{unique_id}_{filename}")
            file.save(input_path)
            
            # Compression logic
            quality = 50 # Default Medium
            if quality_val == 'low': quality = 80
            elif quality_val == 'high': quality = 20
            
            output_filename = f"compressed_{unique_id}_{filename}"
            output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
            
            try:
                img = Image.open(input_path)
                if img.mode in ("RGBA", "P"):
                    img = img.convert("RGB")
                img.save(output_path, optimize=True, quality=quality)
                return render_template('image_compress.html', download_url=output_filename, success=True)
            except Exception as e:
                flash(f'Error processing image: {str(e)}', 'error')
                return redirect(request.url)
        else:
            flash('Invalid file type.', 'error')
            return redirect(request.url)
            
    return render_template('image_compress.html')

@app.route('/merge', methods=['GET', 'POST'])
def pdf_merge():
    if request.method == 'POST':
        cleanup_old_files()
        if 'files' not in request.files:
            flash('No files part', 'error')
            return redirect(request.url)
        
        files = request.files.getlist('files')
        if not files or files[0].filename == '':
            flash('No files selected', 'error')
            return redirect(request.url)
        
        merger = PdfMerger()
        processed_any = False
        
        try:
            for file in files:
                if file and allowed_file(file.filename, ALLOWED_PDF_EXTENSIONS):
                    filename = secure_filename(file.filename)
                    unique_id = str(uuid.uuid4())[:8]
                    path = os.path.join(app.config['UPLOAD_FOLDER'], f"{unique_id}_{filename}")
                    file.save(path)
                    merger.append(path)
                    processed_any = True
            
            if processed_any:
                output_filename = f"merged_{str(uuid.uuid4())[:8]}.pdf"
                output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
                merger.write(output_path)
                merger.close()
                return render_template('pdf_merge.html', download_url=output_filename, success=True)
            else:
                flash('No valid PDF files found.', 'error')
                return redirect(request.url)
        except Exception as e:
            flash(f'Error merging PDFs: {str(e)}', 'error')
            return redirect(request.url)
            
    return render_template('pdf_merge.html')

@app.route('/convert', methods=['GET', 'POST'])
def image_convert():
    if request.method == 'POST':
        cleanup_old_files()
        if 'file' not in request.files:
            flash('No file part', 'error')
            return redirect(request.url)
        
        file = request.files['file']
        target_format = request.form.get('format', 'PNG').upper()
        
        if file.filename == '':
            flash('No selected file', 'error')
            return redirect(request.url)
        
        if file and allowed_file(file.filename, ALLOWED_IMAGE_EXTENSIONS):
            filename = secure_filename(file.filename)
            unique_id = str(uuid.uuid4())[:8]
            input_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{unique_id}_{filename}")
            file.save(input_path)
            
            output_filename = f"converted_{unique_id}.{target_format.lower()}"
            output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
            
            try:
                img = Image.open(input_path)
                if target_format in ['JPG', 'JPEG']:
                    if img.mode in ("RGBA", "P"):
                        img = img.convert("RGB")
                    img.save(output_path, "JPEG")
                else:
                    img.save(output_path, target_format)
                
                return render_template('image_convert.html', download_url=output_filename, success=True)
            except Exception as e:
                flash(f'Error converting image: {str(e)}', 'error')
                return redirect(request.url)
        else:
            flash('Invalid file type.', 'error')
            return redirect(request.url)
            
    return render_template('image_convert.html')

@app.route('/image-to-pdf', methods=['GET', 'POST'])
def image_to_pdf():
    if request.method == 'POST':
        cleanup_old_files()
        files = request.files.getlist('files')
        if not files or files[0].filename == '':
            flash('No files selected', 'error')
            return redirect(request.url)
        
        images = []
        try:
            for file in files:
                if file and allowed_file(file.filename, ALLOWED_IMAGE_EXTENSIONS):
                    filename = secure_filename(file.filename)
                    path = os.path.join(app.config['UPLOAD_FOLDER'], f"{str(uuid.uuid4())[:8]}_{filename}")
                    file.save(path)
                    img = Image.open(path)
                    if img.mode != 'RGB':
                        img = img.convert('RGB')
                    images.append(img)
            
            if images:
                output_filename = f"images_to_pdf_{str(uuid.uuid4())[:8]}.pdf"
                output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
                images[0].save(output_path, save_all=True, append_images=images[1:])
                return render_template('image_to_pdf.html', download_url=output_filename, success=True)
            else:
                flash('No valid images uploaded', 'error')
                return redirect(request.url)
        except Exception as e:
            flash(f'Error: {str(e)}', 'error')
            return redirect(request.url)
            
    return render_template('image_to_pdf.html')

@app.route('/pdf-to-image', methods=['GET', 'POST'])
def pdf_to_image():
    if request.method == 'POST':
        cleanup_old_files()
        if 'file' not in request.files:
            flash('No file part', 'error')
            return redirect(request.url)
        
        file = request.files['file']
        if file.filename == '':
            flash('No file selected', 'error')
            return redirect(request.url)
        
        if file and allowed_file(file.filename, ALLOWED_PDF_EXTENSIONS):
            filename = secure_filename(file.filename)
            unique_id = str(uuid.uuid4())[:8]
            input_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{unique_id}_{filename}")
            file.save(input_path)
            
            try:
                # Try pdf2image (poppler) first
                try:
                    if os.path.exists(POPPLER_PATH):
                        pages = convert_from_path(input_path, poppler_path=POPPLER_PATH)
                    else:
                        pages = convert_from_path(input_path)
                    
                    if len(pages) == 1:
                        output_filename = f"pdf_page_{unique_id}.png"
                        output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
                        pages[0].save(output_path, 'PNG')
                    else:
                        output_filename = f"pdf_pages_{unique_id}.zip"
                        output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
                        with zipfile.ZipFile(output_path, 'w') as zipf:
                            for i, page in enumerate(pages):
                                page_filename = f"page_{i+1}.png"
                                temp_path = os.path.join(app.config['OUTPUT_FOLDER'], f"temp_{unique_id}_{page_filename}")
                                page.save(temp_path, 'PNG')
                                zipf.write(temp_path, page_filename)
                                os.remove(temp_path)
                    
                    return render_template('pdf_to_image.html', download_url=output_filename, success=True)
                
                except Exception:
                    # Fallback to PyMuPDF
                    fallback_images = pdf_to_images_fallback(input_path, app.config['OUTPUT_FOLDER'], unique_id)
                    
                    if len(fallback_images) == 1:
                        # Move fallback file to final output_filename
                        output_filename = f"pdf_page_{unique_id}.png"
                        output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
                        os.rename(fallback_images[0][0], output_path)
                    else:
                        output_filename = f"pdf_pages_{unique_id}.zip"
                        output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
                        with zipfile.ZipFile(output_path, 'w') as zipf:
                            for temp_file, page_name in fallback_images:
                                zipf.write(temp_file, page_name)
                                os.remove(temp_file)
                    
                    return render_template('pdf_to_image.html', download_url=output_filename, success=True)
                    
            except Exception as e:
                flash(f"Error converting PDF: {str(e)}", "error")
                return redirect(request.url)
        else:
            flash('Invalid file type. Only PDF allowed.', 'error')
            return redirect(request.url)
            
    return render_template('pdf_to_image.html')

@app.route('/download/<filename>')
def download_file(filename):
    file_path = os.path.join(app.config['OUTPUT_FOLDER'], filename)
    
    # Basic mimetype detection
    ext = filename.rsplit('.', 1)[-1].lower()
    mimetypes = {
        'pdf': 'application/pdf',
        'png': 'image/png',
        'jpg': 'image/jpeg',
        'jpeg': 'image/jpeg',
        'webp': 'image/webp',
        'zip': 'application/zip',
        'txt': 'text/plain',
        'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
        'xlsx': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
        'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation'
    }
    mimetype = mimetypes.get(ext, 'application/octet-stream')
    
    return send_file(
        file_path,
        as_attachment=True,
        download_name=filename,
        mimetype=mimetype
    )

@app.route('/preview/<filename>')
def preview_file(filename):
    file_path = os.path.join(app.config['OUTPUT_FOLDER'], filename)
    ext = filename.rsplit('.', 1)[-1].lower()
    mimetypes = {
        'pdf': 'application/pdf',
        'png': 'image/png',
        'jpg': 'image/jpeg',
        'jpeg': 'image/jpeg',
        'webp': 'image/webp',
        'txt': 'text/plain',
    }
    return send_file(
        file_path,
        as_attachment=False,
        download_name=filename,
        mimetype=mimetypes.get(ext, 'application/octet-stream')
    )

if __name__ == '__main__':
    app.run(debug=True, port=5000)
